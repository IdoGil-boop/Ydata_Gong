import pandas as pd
import pptx
import re
import numpy as np
import io


def remove_brackets(text):
    print(f"text: {text}")
    while text.startswith(("[", "{")):
        text = text[1:]
    while text.endswith(("]", "}")):
        text = text[:-1]
    print(f"text after: {text}")
    return text


def check_if_text_is_placeholder(text):
    return text.strip().startswith(("[", "{")) and text.strip().endswith(("]", "}"))


def check_if_placeholder_is_inside_text(text):
    """
    Check if text contains any placeholders (text wrapped in [], [[]], or {})

    Args:
        text: String to check for placeholders

    Returns:
        tuple: (bool, list) - True if text contains placeholders and list of matches, False and empty list otherwise
    """
    patterns = [
        r"[\[{]+.*?[\]}]+",  # Matches text surrounded by one or more [ or { at start and one or more ] or } at end
    ]

    all_matches = []
    for pattern in patterns:
        matches = list(re.finditer(pattern, text))
        all_matches.extend(matches)

    return bool(all_matches), all_matches

def replace_placeholder_with_text(text, placeholder_text, replacement_text):
    return text.replace(placeholder_text, replacement_text)


# Function to replace text with replacement text if it's wrapped in [[]], [] or {}
def replace_placeholders(presentation, replacement_texts=None):
    if replacement_texts is None:
        replacement_texts = [["PLACEHOLDER" for _ in range(100)] for _ in range(len(presentation.slides))]
    
    for slide_idx, slide in enumerate(presentation.slides):
        placeholder_count = 0
        
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue

            for paragraph in shape.text_frame.paragraphs:
                paragraph_text = paragraph.text

                # Check if the entire paragraph is a placeholder
                if check_if_text_is_placeholder(paragraph_text):
                    # Get replacement text for this placeholder
                    replacement = "PLACEHOLDER"
                    if slide_idx < len(replacement_texts) and placeholder_count < len(replacement_texts[slide_idx]):
                        replacement = replacement_texts[slide_idx][placeholder_count]
                    placeholder_count += 1
                    
                    # Clear existing runs
                    for idx in range(len(paragraph.runs)):
                        if idx == 0:
                            # Process text to remove brackets/braces from start and end
                            inner_text = remove_brackets(paragraph_text.strip())
                            paragraph.runs[0].text = replacement
                        else:
                            paragraph.runs[idx].text = ""
                else:
                    # Process each run in the paragraph
                    for run in paragraph.runs:
                        run_text = run.text
                        has_placeholders, matches = check_if_placeholder_is_inside_text(
                            run_text
                        )
                        if has_placeholders:
                            # Replace all instances of placeholders with appropriate replacement text
                            modified_text = run_text
                            for match in matches:
                                placeholder_text = match.group()
                                replacement = "PLACEHOLDER"
                                if slide_idx < len(replacement_texts) and placeholder_count < len(replacement_texts[slide_idx]):
                                    replacement = replacement_texts[slide_idx][placeholder_count]
                                placeholder_count += 1
                                
                                modified_text = replace_placeholder_with_text(
                                    modified_text, placeholder_text, replacement
                                )
                            run.text = modified_text


# For visualization, extract the text to show the changes
def extract_formatted_text(presentation):
    all_text_boxes = []

    for slide_index, slide in enumerate(presentation.slides):
        slide_text_boxes = []

        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue

            # Process each paragraph in the text frame
            paragraphs_text = []
            current_text_block = []

            for paragraph in shape.text_frame.paragraphs:
                paragraph_text = paragraph.text.strip()

                # If empty paragraph, might indicate a line break
                if not paragraph_text:
                    # If we have accumulated text, add it as a text box and reset
                    if current_text_block:
                        paragraphs_text.append("\n".join(current_text_block))
                        current_text_block = []
                else:
                    current_text_block.append(paragraph_text)

            # Add any remaining text
            if current_text_block:
                paragraphs_text.append("\n".join(current_text_block))

            # Add all text blocks from this shape
            slide_text_boxes.extend(paragraphs_text)

        all_text_boxes.append(
            {"slide_index": slide_index + 1, "text_boxes": slide_text_boxes}
        )

    return all_text_boxes


# Create a DataFrame to visualize the modified text
def create_df_to_visualize_prs(prs):
    text_df = pd.DataFrame(
        [
            {"Slide": item["slide_index"], "Text Box": i + 1, "Content": content}
            for item in extract_formatted_text(prs)
            for i, content in enumerate(item["text_boxes"])
        ]
    )


def save_prs(path, prs):
    if type(prs) is pptx.presentation.Presentation:
        prs.save(path)
    print(f"prs saved to {path}")
    return path


# change order of slides
def change_order_of_slides(
    prs, path=r"modified_presentations\new_prs.pptx", new_order=None
):
    """
    Reorders slides in a presentation according to the specified order.

    Args:
        prs: The presentation object (pptx.presentation.Presentation)
        new_order: List of slide indices (1-based) in the desired order

    Returns:
        A new presentation with reordered slides
    """
    # save copy of prs
    new_prs = pptx.Presentation(
        save_prs(r"modified_presentations\presentation_copy.pptx", prs)
    )

    # if new_order is not provided, randomize the order
    if new_order is None:
        new_order = np.random.permutation(np.arange(1, len(prs.slides) + 1)).tolist()

    # Check if new_order indices are valid
    max_index = len(prs.slides)
    for idx in new_order:
        if idx < 1 or idx > max_index:
            raise ValueError(f"Slide index {idx} is out of range (1-{max_index})")

    # Make a copy of the slide id elements:
    slide_id_elements = list(new_prs.slides._sldIdLst)

    # Clear them all out of the official list
    for sldId in slide_id_elements:
        new_prs.slides._sldIdLst.remove(sldId)

    # Re-insert in the new order
    for idx in new_order:
        print(f"idx: {idx}")
        new_prs.slides._sldIdLst.append(slide_id_elements[idx - 1])

    # save new prs
    return save_prs(path, new_prs)


def replace_placeholders_in_xml(
    input_path,
    output_path=r"modified_presentations\presentation_with_placeholders.pptx",
    replacement_texts=None
):
    """
    Loads a presentation and modifies the XML directly to replace placeholder text patterns.

    Args:
        input_path: Path to the original presentation
        output_path: Path where modified presentation will be saved
        replacement_texts: 2D array where rows are slide numbers and columns are placeholder indices

    Returns:
        Path to the modified presentation
    """
    import shutil
    from zipfile import ZipFile
    import re
    import xml.etree.ElementTree as ET

    # Create a copy of the presentation
    shutil.copy2(input_path, output_path)

    # Open the presentation as a zip file
    with ZipFile(output_path, "r") as zin:
        # Read all slide XML files
        slide_files = [f for f in zin.namelist() if f.startswith("ppt/slides/slide")]

        for slide_idx, slide_file in enumerate(slide_files):
            # Read the slide XML
            xml_content = zin.read(slide_file)
            root = ET.fromstring(xml_content)

            # Define namespace map
            nsmap = {
                "a": "http://schemas.openxmlformats.org/drawingml/2006/main",
                "p": "http://schemas.openxmlformats.org/presentationml/2006/main",
            }

            # Find all text elements
            text_elements = root.findall(".//a:t", namespaces=nsmap)
            
            placeholder_count = 0
            for elem in text_elements:
                text = elem.text if elem.text else ""

                # Check if entire text is a placeholder
                if check_if_text_is_placeholder(text):
                    # Get replacement text for this placeholder
                    replacement = "PLACEHOLDER"
                    if replacement_texts and slide_idx < len(replacement_texts) and placeholder_count < len(replacement_texts[slide_idx]):
                        replacement = replacement_texts[slide_idx][placeholder_count]
                    placeholder_count += 1
                    
                    # Remove all brackets from the text using remove_brackets function
                    inner_text = remove_brackets(text.strip())
                    elem.text = replacement
                else:
                    # Check for placeholders within text
                    has_placeholders, matches = check_if_placeholder_is_inside_text(
                        text
                    )
                    if has_placeholders:
                        modified_text = text
                        for match in matches:
                            placeholder_text = match.group()
                            
                            # Get replacement text for this placeholder
                            replacement = "PLACEHOLDER"
                            if replacement_texts and slide_idx < len(replacement_texts) and placeholder_count < len(replacement_texts[slide_idx]):
                                replacement = replacement_texts[slide_idx][placeholder_count]
                            placeholder_count += 1
                            
                            modified_text = replace_placeholder_with_text(
                                modified_text, placeholder_text, replacement
                            )
                        elem.text = modified_text

            # Write modified XML back to the file
            with ZipFile(output_path, "a") as zout:
                zout.writestr(slide_file, ET.tostring(root, encoding="UTF-8"))

    return output_path


# Example usage
if __name__ == "__main__":
    import pptx

    # Load the presentation
    prs = pptx.Presentation(
        r"/Users/ido.gil-ext@gong.io/Ydata_Gong/original_presentations/Gong Sample Deck Slides with Placeholders Template.pptx"
    )

    # Example of replacement texts - 2D array where rows are slides and columns are placeholders
    replacement_texts = [
        ["Title 1", "Subtitle 1"],  # Replacements for slide 1
        ["Point 1", "Point 2", "Point 3"],  # Replacements for slide 2
        ["Conclusion Text"]  # Replacements for slide 3
    ]

    # Replace placeholders in the XML
    replace_placeholders_in_xml(
        r"/Users/ido.gil-ext@gong.io/Ydata_Gong/original_presentations/Gong Sample Deck Slides with Placeholders Template.pptx",
        r"modified_presentations/presentation_with_placeholders.pptx",
        replacement_texts
    )

    # # Save the reordered presentation
    # save_prs(
    #     r'modified_presentations\presentation_with_reordered_slides.pptx',
    #     reordered_prs
    # )
