""" run with streamlit run placeholder_replace_app.py """


import uuid
import xml.etree.ElementTree as ET
from zipfile import ZipFile, ZIP_DEFLATED  # Import ZIP_DEFLATED for compression
import shutil
import tempfile
import os
import io
import re
import numpy as np
import pandas as pd
import streamlit as st

try:
    import pptx
except ImportError:
    st.error("Could not import pptx. Try installing with: pip install python-pptx")

# --- Paste your existing helper functions here (remove_brackets, etc.) ---


def remove_brackets(text):
    # print(f"text: {text}") # Commented out for cleaner Streamlit output
    if not isinstance(text, str):
        return text
    temp_text = text.strip()
    # Make it more robust for single/double brackets/braces
    while temp_text.startswith(("[", "{")) or temp_text.startswith(("[[", "{{")):
        temp_text = temp_text[1:] if not (
            temp_text.startswith(("[[", "{{"))) else temp_text[2:]
    while temp_text.endswith(("]", "}")) or temp_text.endswith(("]]", "}}")):
        temp_text = temp_text[:-1] if not (
            temp_text.endswith(("]]", "}}"))) else temp_text[:-2]
    # print(f"text after: {temp_text}") # Commented out
    return temp_text.strip()  # Add strip at the end too


def check_if_text_is_placeholder(text):
    if not isinstance(text, str):
        return False
    # Allow single/double brackets/braces
    return text.strip().startswith(("[", "{")) and text.strip().endswith(("]", "}"))


def check_if_placeholder_is_inside_text(text):
    if not isinstance(text, str):
        return False, []
    # Pattern to find placeholders like [text], {text}, [[text]], {{text}} non-greedily
    pattern = r"[\[{]{1,2}.*?[\]}]{1,2}"
    all_matches = []
    try:
        # Use re.findall to get non-overlapping matches directly
        matches = re.findall(pattern, text)
        # Convert matches to match objects if needed, or just use the strings
        # For consistency with the old code structure returning match objects:
        all_matches = list(re.finditer(pattern, text))
    except re.error:
        pass  # Ignore regex errors on potentially malformed strings
    return bool(all_matches), all_matches


def replace_placeholder_with_text(text, placeholder_text, replacement_text):
    if not isinstance(text, str):
        return text
    # Ensure replacement is a string
    if replacement_text is None:
        replacement_text = ""
    elif not isinstance(replacement_text, str):
        replacement_text = str(replacement_text)

    try:
        # Use regex for replacement to handle potential special characters in placeholder_text
        # Escape the placeholder text to treat it literally in the regex
        escaped_placeholder = re.escape(placeholder_text)
        # Using count=1 ensures we only replace the first occurrence if duplicates exist
        # within the same text run and multiple matches were found (though the calling logic
        # should ideally handle this by iterating matches)
        return re.sub(escaped_placeholder, replacement_text, text, count=1)
        # If exact string matching is sufficient and faster:
        # return text.replace(placeholder_text, replacement_text, 1) # replace only first instance per call

    except (TypeError, re.error):
        # Handle cases where text might not be a string or regex fails
        return text

# --- REVISED replace_placeholders_in_xml Function ---


def replace_placeholders_in_xml(
    input_path,
    output_path,
    replacement_texts=None
):
    """
    Loads a presentation, modifies the XML for specified slides, and saves
    to a new presentation file safely.

    Args:
        input_path: Path to the original presentation.
        output_path: Path where the modified presentation will be saved.
        replacement_texts: 2D list where outer index = slide index,
                         inner list contains replacements for that slide in order.

    Returns:
        Path to the modified presentation or None if an error occurs.
    """
    if replacement_texts is None:
        replacement_texts = []

    # Use a temporary file for the new zip archive to avoid conflicts
    temp_zip_file = None
    try:
        temp_dir = os.path.dirname(output_path)
        os.makedirs(temp_dir, exist_ok=True)  # Ensure output directory exists

        # Create a named temporary file for the output ZIP
        temp_fd, temp_zip_path = tempfile.mkstemp(suffix=".pptx", dir=temp_dir)
        os.close(temp_fd)  # Close the file descriptor, we only need the path
        temp_zip_file = temp_zip_path  # Store path for cleanup

        nsmap = {
            "a": "http://schemas.openxmlformats.org/drawingml/2006/main",
            "p": "http://schemas.openxmlformats.org/presentationml/2006/main",
            "r": "http://schemas.openxmlformats.org/officeDocument/2006/relationships"
        }
        modified_slide_data = {}  # Store modified XML content keyed by slide file name

        # --- Pass 1: Read and modify slide XMLs in memory ---
        with ZipFile(input_path, "r") as zin:
            slide_files = sorted(
                [f for f in zin.namelist() if f.startswith(
                    "ppt/slides/slide") and f.endswith(".xml")],
                key=lambda x: int(re.search(r'(\d+)', x).group(1))
            )

            for slide_idx, slide_file in enumerate(slide_files):
                try:
                    xml_content = zin.read(slide_file)
                    root = ET.fromstring(xml_content)
                    text_elements = root.findall(".//a:t", namespaces=nsmap)
                    placeholder_count_on_slide = 0
                    slide_was_modified = False  # Track if modifications were made

                    for elem in text_elements:
                        original_text = elem.text if elem.text else ""
                        current_elem_text = original_text  # Work on a copy

                        # 1. Check if the entire text element is a placeholder
                        if check_if_text_is_placeholder(current_elem_text):
                            if (replacement_texts and
                                slide_idx < len(replacement_texts) and
                                    placeholder_count_on_slide < len(replacement_texts[slide_idx])):
                                replacement = replacement_texts[slide_idx][placeholder_count_on_slide]
                                elem.text = replacement if replacement is not None else ""
                                placeholder_count_on_slide += 1
                                slide_was_modified = True
                            else:
                                # Optional: Handle case where replacement is missing
                                # elem.text = "[MISSING REPLACEMENT]"
                                pass  # Or leave as is, or clear it: elem.text = ""

                        # 2. Check for placeholders *within* the text element
                        else:
                            has_placeholders, matches = check_if_placeholder_is_inside_text(
                                current_elem_text)
                            if has_placeholders:
                                # Store changes temporarily to avoid modifying text while iterating matches
                                temp_modified_text = current_elem_text
                                for match in matches:
                                    placeholder_text = match.group(0)
                                    if (replacement_texts and
                                        slide_idx < len(replacement_texts) and
                                            placeholder_count_on_slide < len(replacement_texts[slide_idx])):

                                        replacement = replacement_texts[slide_idx][placeholder_count_on_slide]
                                        # Use the replacement function to substitute *this* match
                                        # Important: Pass temp_modified_text to replace within the potentially already modified string
                                        temp_modified_text = replace_placeholder_with_text(
                                            temp_modified_text, placeholder_text, replacement
                                        )
                                        placeholder_count_on_slide += 1
                                        slide_was_modified = True
                                    else:
                                        # Optional: Handle missing replacement for inner placeholder
                                        # temp_modified_text = replace_placeholder_with_text(
                                        #    temp_modified_text, placeholder_text, "[MISSING]"
                                        # )
                                        pass

                                # Update the element text only once after processing all matches for this element
                                elem.text = temp_modified_text

                    # If the slide was modified, store the new XML content
                    if slide_was_modified:
                        # Important: Use 'unicode' encoding for ET.tostring to get a string,
                        # then encode to UTF-8 bytes for writing to zip.
                        # Adding xml_declaration is crucial for PowerPoint.
                        modified_xml_string = ET.tostring(
                            root, encoding='unicode', xml_declaration=True)
                        modified_slide_data[slide_file] = modified_xml_string.encode(
                            'utf-8')

                except ET.ParseError as e:
                    st.warning(
                        f"Skipping slide {slide_idx + 1} ({slide_file}) due to XML parsing error: {e}")
                    continue
                except Exception as e:
                    st.error(
                        f"An unexpected error occurred processing slide {slide_idx + 1} ({slide_file}): {e}")
                    continue  # Skip to next slide

        # --- Pass 2: Create the new ZIP file ---
        with ZipFile(temp_zip_path, "w", compression=ZIP_DEFLATED) as zout:
            with ZipFile(input_path, "r") as zin:
                for item in zin.infolist():
                    # Read the content from the original zip
                    file_content = zin.read(item.filename)

                    # If it's a modified slide, write the modified content
                    if item.filename in modified_slide_data:
                        # Make sure not to write empty data if modification failed silently
                        if modified_slide_data[item.filename]:
                            zout.writestr(
                                item, modified_slide_data[item.filename])
                        else:  # Fallback to original content if modified data is bad
                            st.warning(
                                f"Modified data for {item.filename} was empty, writing original.")
                            zout.writestr(item, file_content)
                    # Otherwise, write the original content
                    else:
                        zout.writestr(item, file_content)

        # --- Final Step: Replace the target output path with the temp file ---
        # shutil.move is generally robust for this replacement.
        shutil.move(temp_zip_path, output_path)
        temp_zip_file = None  # Prevent cleanup since move was successful

        return output_path

    except FileNotFoundError:
        st.error(f"Error: Input file not found at {input_path}")
        return None
    except ET.ParseError as e:
        st.error(
            f"Fatal XML Parsing Error: {e}. The input file might be corrupt or not a valid PPTX.")
        return None
    except Exception as e:
        st.error(
            f"An unexpected error occurred during placeholder replacement: {e}")
        import traceback
        # More detailed error for debugging
        st.error(f"Traceback: {traceback.format_exc()}")
        return None
    finally:
        # --- Cleanup: Ensure temporary zip file is deleted if an error occurred ---
        if temp_zip_file and os.path.exists(temp_zip_file):
            try:
                os.remove(temp_zip_file)
            except OSError as e:
                st.warning(
                    f"Could not remove temporary file {temp_zip_file}: {e}")


# --- Your extract_placeholders_for_ui function (should be okay, but double-check indexing logic matches the revised replacement function) ---
# Make sure the placeholder_count_on_slide increments identically to the replacement function.
# If you encounter issues where replacements are still wrong, add print statements inside both
# loops (in extract and replace) to compare the order and indices assigned to placeholders.
def extract_placeholders_for_ui(input_path):
    """
    Scans the presentation XML to find placeholders for UI display.
    Ensure placeholder counting logic mirrors the replacement function.
    """
    placeholders_found = []
    try:
        with ZipFile(input_path, "r") as zin:
            nsmap = {
                "a": "http://schemas.openxmlformats.org/drawingml/2006/main",
                "p": "http://schemas.openxmlformats.org/presentationml/2006/main",
            }

            slide_files = sorted([f for f in zin.namelist() if f.startswith("ppt/slides/slide") and f.endswith(".xml")],
                                 key=lambda x: int(re.search(r'(\d+)', x).group(1)))

            for slide_idx, slide_file in enumerate(slide_files):
                try:
                    xml_content = zin.read(slide_file)
                    root = ET.fromstring(xml_content)
                    text_elements = root.findall(".//a:t", namespaces=nsmap)

                    placeholder_count_on_slide = 0
                    for elem_idx, elem in enumerate(text_elements):
                        text = elem.text if elem.text else ""

                        # Check if the entire text is a placeholder
                        if check_if_text_is_placeholder(text):
                            placeholders_found.append({
                                'slide_index': slide_idx,
                                'placeholder_index_on_slide': placeholder_count_on_slide,
                                'original_text': text,
                                'unique_key': f"ph_{slide_idx}_{placeholder_count_on_slide}_{uuid.uuid4()}"
                            })
                            placeholder_count_on_slide += 1  # Increment counter
                        else:
                            # Check for placeholders within the text
                            has_placeholders, matches = check_if_placeholder_is_inside_text(
                                text)
                            if has_placeholders:
                                # IMPORTANT: Increment the counter for EACH match found within the text
                                for match in matches:
                                    placeholder_text = match.group(0)
                                    placeholders_found.append({
                                        'slide_index': slide_idx,
                                        # The index corresponds to this specific placeholder's position in the sequence for this slide
                                        'placeholder_index_on_slide': placeholder_count_on_slide,
                                        'original_text': placeholder_text,
                                        'unique_key': f"ph_{slide_idx}_{placeholder_count_on_slide}_{uuid.uuid4()}"
                                    })
                                    placeholder_count_on_slide += 1  # Increment counter for each placeholder found
                except ET.ParseError as e:
                    st.warning(
                        f"Skipping slide {slide_idx + 1} ({slide_file}) during placeholder extraction due to XML parsing error: {e}")
                    continue
                except Exception as e:
                    st.warning(
                        f"Error extracting placeholders from slide {slide_idx + 1} ({slide_file}): {e}")
                    continue

        # Optional: Sort placeholders based on index before returning, although the Streamlit part sorts later.
        # placeholders_found.sort(key=lambda x: (x['slide_index'], x['placeholder_index_on_slide']))
        return placeholders_found

    except FileNotFoundError:
        st.error(
            f"Error: Input file not found at {input_path} during extraction.")
        return []
    except Exception as e:
        st.error(
            f"An unexpected error occurred during placeholder extraction: {e}")
        return []


st.set_page_config(layout="wide")
st.title("PowerPoint Placeholder Editor")

# --- Session State Initialization ---
if 'uploaded_file_info' not in st.session_state:
    st.session_state.uploaded_file_info = None
if 'placeholders' not in st.session_state:
    st.session_state.placeholders = []
if 'modified_file_path' not in st.session_state:
    st.session_state.modified_file_path = None
if 'processing_done' not in st.session_state:
    st.session_state.processing_done = False
if 'run_key' not in st.session_state:  # To force reset when new file uploaded
    st.session_state.run_key = str(uuid.uuid4())
if 'error_message' not in st.session_state:  # To display errors prominently
    st.session_state.error_message = None


# --- File Upload ---
uploaded_file = st.file_uploader(
    "Choose a PowerPoint file (.pptx)", type="pptx", key=st.session_state.run_key)

# Display persistent error messages if any
if st.session_state.error_message:
    st.error(st.session_state.error_message)
    st.session_state.error_message = None  # Clear after displaying

if uploaded_file is not None:
    if st.session_state.uploaded_file_info is None or st.session_state.uploaded_file_info['name'] != uploaded_file.name:
        st.session_state.placeholders = []
        st.session_state.modified_file_path = None
        st.session_state.processing_done = False
        st.session_state.error_message = None  # Clear previous errors

        with tempfile.NamedTemporaryFile(delete=False, suffix=".pptx") as tmp_file:
            tmp_file.write(uploaded_file.getvalue())
            temp_input_path = tmp_file.name

        st.session_state.uploaded_file_info = {
            'name': uploaded_file.name, 'path': temp_input_path}
        st.success(
            f"Uploaded '{uploaded_file.name}'. Extracting placeholders...")

        with st.spinner("Scanning presentation for placeholders..."):
            st.session_state.placeholders = extract_placeholders_for_ui(
                temp_input_path)

        if not st.session_state.placeholders:
            st.warning(
                "No placeholders (like `[text]` or `{text}`) found in the presentation.")
        else:
            st.info(
                f"Found {len(st.session_state.placeholders)} placeholders. Enter replacements below.")

    if st.session_state.placeholders:
        st.subheader("Edit Placeholders")
        placeholders_by_slide = {}
        for ph in st.session_state.placeholders:
            slide_num = ph['slide_index'] + 1
            if slide_num not in placeholders_by_slide:
                placeholders_by_slide[slide_num] = []
            # Ensure placeholders within a slide are ordered correctly for display
            placeholders_by_slide[slide_num].append(ph)

        # Sort placeholders within each slide list by their on-slide index
        for slide_num in placeholders_by_slide:
            placeholders_by_slide[slide_num].sort(
                key=lambda x: x['placeholder_index_on_slide'])

        with st.form(key='placeholder_form'):
            replacement_inputs = {}
            for slide_num in sorted(placeholders_by_slide.keys()):
                st.markdown(f"**Slide {slide_num}**")
                cols = st.columns(2)
                # Display headers once per slide
                cols[0].write("**Original Placeholder:**")
                cols[1].write("**Replacement Text:**")

                # Keep track of inputs created for this slide
                inputs_this_slide = []
                for ph in placeholders_by_slide[slide_num]:
                    # Use markdown in the first column for the placeholder text
                    cols[0].markdown(f"`{ph['original_text']}`")
                    # Use the second column for the text input
                    # Make label more informative but keep it hidden
                    label = f"Input for '{ph['original_text']}' (Slide {slide_num}, Index {ph['placeholder_index_on_slide']})"
                    # Use the unique key for the widget
                    default_value = remove_brackets(
                        ph['original_text'])  # Pre-fill suggestion
                    widget_key = ph['unique_key']
                    replacement_inputs[widget_key] = cols[1].text_input(
                        label=label,
                        # Try to preserve input on re-run
                        value=st.session_state.get(widget_key, default_value),
                        key=widget_key,
                        label_visibility="collapsed"
                    )
                    inputs_this_slide.append(widget_key)  # Track keys used

                st.divider()

            submitted = st.form_submit_button(
                "Apply Changes and Generate Modified Presentation")

            if submitted:
                st.session_state.processing_done = False
                st.session_state.modified_file_path = None
                st.session_state.error_message = None

                # --- Prepare replacement_texts structure ---
                max_slide_idx = -1
                if st.session_state.placeholders:
                    max_slide_idx = max(p['slide_index']
                                        for p in st.session_state.placeholders)

                replacement_texts_structured = [[]
                                                for _ in range(max_slide_idx + 1)]

                # Crucially, sort the GLOBAL list of placeholders exactly as they were found/indexed
                sorted_placeholders = sorted(st.session_state.placeholders, key=lambda x: (
                    x['slide_index'], x['placeholder_index_on_slide']))

                # Populate the structure based on user inputs from the form
                successful_population = True
                for ph in sorted_placeholders:
                    slide_idx = ph['slide_index']
                    ph_key = ph['unique_key']
                    # Get user input using the unique key
                    user_input = replacement_inputs.get(ph_key, "")

                    if slide_idx < len(replacement_texts_structured):
                        replacement_texts_structured[slide_idx].append(
                            user_input)
                    else:
                        st.error(
                            f"Structure error processing placeholder for slide index {slide_idx}. Max index was {max_slide_idx}.")
                        st.session_state.error_message = "Internal error structuring replacements. Please try again."
                        successful_population = False
                        break  # Stop processing if structure is wrong

                if successful_population:
                    st.info("Applying replacements... Please wait.")
                    input_pptx_path = st.session_state.uploaded_file_info['path']
                    temp_output_dir = tempfile.mkdtemp()
                    # Ensure the output filename is distinct, maybe add a timestamp or UUID
                    output_filename = f"modified_{uuid.uuid4().hex[:8]}_{st.session_state.uploaded_file_info['name']}"
                    output_pptx_path = os.path.join(
                        temp_output_dir, output_filename)

                    # --- Call the revised replacement function ---
                    modified_path = replace_placeholders_in_xml(
                        input_path=input_pptx_path,
                        output_path=output_pptx_path,
                        replacement_texts=replacement_texts_structured
                    )

                    if modified_path and os.path.exists(modified_path):
                        st.session_state.modified_file_path = modified_path
                        st.session_state.processing_done = True
                        # Use st.rerun() to refresh the state and show the download button immediately
                        st.rerun()
                    else:
                        st.error(
                            "Failed to apply replacements. Check error messages above.")
                        st.session_state.error_message = "Failed to create modified presentation."
                        # Clean up temp output dir if failed
                        if os.path.exists(temp_output_dir):
                            shutil.rmtree(temp_output_dir, ignore_errors=True)

    # --- Download Button ---
    # This block should be outside the `if st.session_state.placeholders:` block
    # but inside the `if uploaded_file is not None:` block
    if st.session_state.processing_done and st.session_state.modified_file_path:
        st.subheader("Download Modified Presentation")
        try:
            # Extract only the filename for the download attribute
            download_filename = os.path.basename(
                st.session_state.modified_file_path)
            with open(st.session_state.modified_file_path, "rb") as fp:
                st.download_button(
                    label="Download Modified PPTX",
                    data=fp,
                    file_name=download_filename,  # Use the generated name
                    mime="application/vnd.openxmlformats-officedocument.presentationml.presentation",
                )
        except FileNotFoundError:
            st.error("Modified file not found. Please try applying changes again.")
            st.session_state.error_message = "Could not find the modified file for download."
            st.session_state.processing_done = False  # Reset state
            st.session_state.modified_file_path = None
        except Exception as e:
            st.error(f"Error preparing download: {e}")
            st.session_state.error_message = f"Error during download preparation: {e}"


else:
    # --- Reset and Cleanup when no file is uploaded ---
    if st.session_state.uploaded_file_info is not None:
        # Clean up temp input file
        if 'path' in st.session_state.uploaded_file_info and st.session_state.uploaded_file_info['path'] and os.path.exists(st.session_state.uploaded_file_info['path']):
            try:
                os.remove(st.session_state.uploaded_file_info['path'])
            except OSError as e:
                st.warning(f"Could not remove temp input file: {e}")
        # Clean up temp output directory if a path exists
        if st.session_state.modified_file_path:
            temp_output_dir = os.path.dirname(
                st.session_state.modified_file_path)
            if os.path.exists(temp_output_dir):
                shutil.rmtree(temp_output_dir, ignore_errors=True)

    # Reset session state variables
    st.session_state.uploaded_file_info = None
    st.session_state.placeholders = []
    st.session_state.modified_file_path = None
    st.session_state.processing_done = False
    st.session_state.error_message = None  # Clear any lingering errors
    st.info("Please upload a PowerPoint (.pptx) file to begin.")
